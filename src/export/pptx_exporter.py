"""
PowerPoint Export Module for Dashboard Reports.
Generates professional PPTX presentations with charts, tables, and analytics.
"""
import os
from datetime import datetime
from typing import Dict, List, Any, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import pandas as pd
from io import BytesIO

from src.utils.logger import get_logger

logger = get_logger(__name__)


class PowerPointExporter:
    """Export dashboard data to PowerPoint presentation."""
    
    # Corporate colors
    COLORS = {
        'primary': RGBColor(30, 136, 229),    # #1E88E5
        'secondary': RGBColor(67, 160, 71),   # #43A047
        'accent': RGBColor(251, 140, 0),      # #FB8C00
        'danger': RGBColor(229, 57, 53),      # #E53935
        'dark': RGBColor(33, 33, 33),
        'light': RGBColor(245, 245, 245),
        'white': RGBColor(255, 255, 255)
    }
    
    def __init__(self, template_path: Optional[str] = None):
        self.template_path = template_path
        self.prs = Presentation(template_path) if template_path else Presentation()
        self._setup_master()
        
    def _setup_master(self):
        """Configure slide master styles."""
        # Set default font
        slide_master = self.prs.slide_master
        text_frame = slide_master.placeholders[0].text_frame if slide_master.placeholders else None
        if text_frame:
            paragraph = text_frame.paragraphs[0]
            font = paragraph.font
            font.name = 'Calibri'
            font.size = Pt(18)
    
    def _add_title_slide(self, title: str, subtitle: str = "") -> None:
        """Add a title slide."""
        slide_layout = self.prs.slide_layouts[0]  # Title Slide
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        title_tf = title_placeholder.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = self.COLORS['primary']
        
        # Set subtitle
        if subtitle and len(slide.placeholders) > 1:
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle
            subtitle_tf = subtitle_placeholder.text_frame
            subtitle_p = subtitle_tf.paragraphs[0]
            subtitle_p.font.size = Pt(24)
            subtitle_p.font.color.rgb = self.COLORS['dark']
    
    def _add_content_slide(self, title: str) -> Any:
        """Add a content slide with title."""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        title_tf = title_placeholder.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = self.COLORS['primary']
        
        return slide
    
    def _add_kpi_slide(self, metrics: Dict[str, Any]) -> None:
        """Add a slide with KPI cards."""
        slide = self._add_content_slide("Ключевые показатели")
        
        # Calculate layout
        kpi_items = list(metrics.items())[:6]  # Max 6 KPIs
        cols = 3 if len(kpi_items) > 3 else len(kpi_items)
        rows = (len(kpi_items) + cols - 1) // cols
        
        card_width = Inches(3.5)
        card_height = Inches(2)
        spacing = Inches(0.3)
        start_x = Inches(0.5)
        start_y = Inches(2)
        
        for idx, (key, value) in enumerate(kpi_items):
            col = idx % cols
            row = idx // cols
            
            x = start_x + col * (card_width + spacing)
            y = start_y + row * (card_height + spacing)
            
            # Create card shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_width, card_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.COLORS['light']
            shape.line.color.rgb = self.COLORS['primary']
            shape.line.width = Pt(2)
            
            # Add metric name
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = key
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['dark']
            p.alignment = PP_ALIGN.CENTER
            
            # Add metric value
            if isinstance(value, dict):
                val = value.get('value', 'N/A')
                change = value.get('change', '0%')
            else:
                val = value
                change = None
            
            p = tf.add_paragraph()
            p.text = str(val)
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['primary']
            p.alignment = PP_ALIGN.CENTER
            
            # Add change indicator
            if change:
                p = tf.add_paragraph()
                change_str = f"{change}" if str(change).startswith('+') or str(change).startswith('-') else f"+{change}"
                p.text = change_str
                p.font.size = Pt(14)
                color = self.COLORS['secondary'] if str(change).startswith('+') else self.COLORS['danger']
                p.font.color.rgb = color
                p.alignment = PP_ALIGN.CENTER
    
    def _add_chart_slide(
        self,
        title: str,
        chart_type: str,
        data: Dict[str, Any],
        position: Tuple[float, float, float, float] = (Inches(1), Inches(2.5), Inches(8), Inches(4.5))
    ) -> None:
        """Add a slide with a chart."""
        slide = self._add_content_slide(title)
        
        x, y, cx, cy = position
        
        if chart_type == 'line':
            chart_data = XyChartData()
            series_data = chart_data.add_series(data.get('series_name', 'Данные'))
            
            x_values = data.get('x', [])
            y_values = data.get('y', [])
            
            for x_val, y_val in zip(x_values, y_values):
                series_data.add_data_point(x_val, y_val)
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.XY_SCATTER_LINES,
                x, y, cx, cy, chart_data
            )
            
        elif chart_type == 'bar':
            chart_data = CategoryChartData()
            chart_data.categories = data.get('categories', [])
            chart_data.add_series(data.get('series_name', 'Значения'), data.get('values', []))
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                x, y, cx, cy, chart_data
            )
            
        elif chart_type == 'pie':
            chart_data = CategoryChartData()
            chart_data.categories = data.get('labels', [])
            chart_data.add_series(data.get('series_name', 'Доля'), data.get('sizes', []))
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                x, y, cx, cy, chart_data
            )
            
            # Format pie chart
            plot = chart.chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_percentage = True
            data_labels.position = XL_LEGEND_POSITION.OUTSIDE_END
        
        # Style the chart
        chart.chart.has_legend = True
        chart.chart.legend.include_in_layout = False
        chart.chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    def _add_table_slide(
        self,
        title: str,
        df: pd.DataFrame,
        max_rows: int = 15
    ) -> None:
        """Add a slide with a data table."""
        slide = self._add_content_slide(title)
        
        # Limit rows
        if len(df) > max_rows:
            df = df.head(max_rows)
        
        # Table dimensions
        rows = len(df) + 1  # +1 for header
        cols = len(df.columns)
        
        x = Inches(0.5)
        y = Inches(2)
        width = Inches(9)
        height = Inches(0.4)
        
        # Adjust height for content
        table_height = min(Inches(5), height * rows)
        
        table = slide.shapes.add_table(rows, cols, x, y, width, table_height).table
        
        # Set column widths
        col_width = width / cols
        for i in range(cols):
            table.columns[i].width = col_width
        
        # Fill header
        for i, col_name in enumerate(df.columns):
            cell = table.cell(0, i)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.COLORS['primary']
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.bold = True
                paragraph.font.color.rgb = self.COLORS['white']
                paragraph.alignment = PP_ALIGN.CENTER
        
        # Fill data
        for row_idx, row in df.iterrows():
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)
                
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.alignment = PP_ALIGN.CENTER
                
                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.COLORS['light']
        
        # Add note if truncated
        if len(df) >= max_rows:
            note = slide.shapes.add_textbox(
                Inches(0.5), y + table_height + Inches(0.2),
                Inches(9), Inches(0.5)
            )
            note_tf = note.text_frame
            note_p = note_tf.paragraphs[0]
            note_p.text = f"Показано первые {max_rows} строк из {len(df)}"
            note_p.font.size = Pt(10)
            note_p.font.italic = True
            note_p.font.color.rgb = self.COLORS['dark']
    
    def _add_summary_slide(self, summary: str) -> None:
        """Add a summary/conclusion slide."""
        slide = self._add_content_slide("Выводы и рекомендации")
        
        # Add text box with summary
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        paragraphs = summary.split('\n')
        for i, para_text in enumerate(paragraphs):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {para_text}" if para_text.strip() else ""
            p.font.size = Pt(16)
            p.space_after = Pt(10)
            p.font.color.rgb = self.COLORS['dark']
    
    def export_dashboard(
        self,
        metrics: Dict[str, Any],
        charts: List[Dict[str, Any]],
        dataframes: Dict[str, pd.DataFrame],
        title: str = "Аналитический отчет",
        subtitle: str = "",
        summary: str = "",
        output_path: Optional[str] = None
    ) -> str:
        """
        Export complete dashboard to PowerPoint.
        
        Args:
            metrics: Dictionary of KPI metrics
            charts: List of chart configurations
            dataframes: Dictionary of DataFrames to include
            title: Presentation title
            subtitle: Presentation subtitle
            summary: Summary/conclusion text
            output_path: Output file path
            
        Returns:
            Path to generated PPTX
        """
        # Reset presentation
        self.prs = Presentation() if not self.template_path else Presentation(self.template_path)
        self._setup_master()
        
        # 1. Title slide
        self._add_title_slide(title, subtitle)
        
        # 2. KPI slide
        if metrics:
            self._add_kpi_slide(metrics)
        
        # 3. Chart slides
        for chart_config in charts:
            chart_type = chart_config.get('type', 'bar')
            chart_title = chart_config.get('title', 'График')
            self._add_chart_slide(chart_title, chart_type, chart_config)
        
        # 4. Data table slides
        for table_name, df in dataframes.items():
            self._add_table_slide(table_name, df)
        
        # 5. Summary slide
        if summary:
            self._add_summary_slide(summary)
        
        # Save presentation
        output = output_path or f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        self.prs.save(output)
        
        logger.info(f"PowerPoint presentation generated: {output}")
        return output
    
    def export_simple_presentation(
        self,
        title: str,
        slides: List[Dict[str, Any]],
        output_path: Optional[str] = None
    ) -> str:
        """Export a simple custom presentation."""
        self.prs = Presentation()
        self._setup_master()
        
        # Title slide
        self._add_title_slide(title)
        
        # Content slides
        for slide_data in slides:
            slide_type = slide_data.get('type', 'content')
            
            if slide_type == 'kpi':
                self._add_kpi_slide(slide_data.get('metrics', {}))
            elif slide_type == 'chart':
                self._add_chart_slide(
                    slide_data.get('title', ''),
                    slide_data.get('chart_type', 'bar'),
                    slide_data.get('data', {})
                )
            elif slide_type == 'table':
                self._add_table_slide(
                    slide_data.get('title', ''),
                    slide_data.get('dataframe', pd.DataFrame())
                )
            elif slide_type == 'summary':
                self._add_summary_slide(slide_data.get('text', ''))
        
        output = output_path or f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        self.prs.save(output)
        
        logger.info(f"Simple PowerPoint presentation generated: {output}")
        return output


def export_to_pptx(
    metrics: Dict[str, Any],
    charts: Optional[List[Dict[str, Any]]] = None,
    dataframes: Optional[Dict[str, pd.DataFrame]] = None,
    output_path: Optional[str] = None,
    title: str = "Аналитический отчет",
    subtitle: str = "",
    summary: str = ""
) -> str:
    """
    Convenience function to export dashboard to PowerPoint.
    
    Args:
        metrics: KPI metrics dictionary
        charts: List of chart configurations
        dataframes: Dictionary of DataFrames
        output_path: Output file path
        title: Presentation title
        subtitle: Presentation subtitle
        summary: Summary text
        
    Returns:
        Path to generated PPTX
    """
    exporter = PowerPointExporter()
    
    charts = charts or []
    dataframes = dataframes or {}
    
    return exporter.export_dashboard(
        metrics=metrics,
        charts=charts,
        dataframes=dataframes,
        output_path=output_path,
        title=title,
        subtitle=subtitle,
        summary=summary
    )
